from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Brand colors
DBC_GREEN = RGBColor(0x32, 0xCD, 0x32)       # Lime green
DBC_DARK_GREEN = RGBColor(0x22, 0x8B, 0x22)  # Darker lime for bg slides
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MED_GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_BG = RGBColor(0xF5, 0xF5, 0xF5)
FONT = "Arial"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SLIDE_W = Inches(13.333)


def add_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18, color=DARK_GRAY, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = FONT
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=DARK_GRAY):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = FONT
        p.space_after = Pt(10)
    return txBox


# ── SLIDE 1: Title ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DBC_DARK_GREEN)

add_text_box(slide, Inches(1.5), Inches(2.0), Inches(10), Inches(1.5),
             "DOLLAR BAR CLUB", 54, WHITE, True, PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(3.5), Inches(10), Inches(0.8),
             "Bar Partner Overview", 28, WHITE, False, PP_ALIGN.CENTER)

add_shape(slide, Inches(5.5), Inches(4.5), Inches(2.3), Pt(2), WHITE)

add_text_box(slide, Inches(1.5), Inches(4.9), Inches(10), Inches(0.6),
             "Powered by BarGlance  |  Austin Pilot Program", 18, WHITE, False, PP_ALIGN.CENTER)


# ── SLIDE 2: What It Is + How It Works (combined) ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.15), DBC_GREEN)

add_text_box(slide, Inches(1), Inches(0.6), Inches(11), Inches(0.8),
             "What is Dollar Bar Club?", 36, DBC_GREEN, True)

add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(0.8),
             "A free digital membership that drives new customers to your bar through exclusive, one-time drink offers. No app download needed - it's a mobile web app.",
             20, DARK_GRAY)

# 4-step flow
steps = [
    ("1", "SIGN UP", "Free membership via\nemail + Austin ZIP."),
    ("2", "DISCOVER", "Members browse bars\nand see your offer."),
    ("3", "VISIT", "GPS confirms they're\nat your venue."),
    ("4", "REDEEM", "Show screen to bartender.\nEnjoy the drink."),
]

for i, (num, title, desc) in enumerate(steps):
    x = Inches(0.8 + i * 3.1)
    y = Inches(3.2)

    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(1.05), y, Inches(0.7), Inches(0.7))
    circle.fill.solid()
    circle.fill.fore_color.rgb = DBC_GREEN
    circle.line.fill.background()
    tf = circle.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(22)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = FONT
    p.alignment = PP_ALIGN.CENTER

    add_text_box(slide, x, y + Inches(0.9), Inches(2.8), Inches(0.5),
                 title, 20, DBC_GREEN, True, PP_ALIGN.CENTER)
    add_text_box(slide, x, y + Inches(1.5), Inches(2.8), Inches(1.2),
                 desc, 15, MED_GRAY, False, PP_ALIGN.CENTER)


# ── SLIDE 3: Why Partner With Us ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.15), DBC_GREEN)

add_text_box(slide, Inches(1), Inches(0.6), Inches(11), Inches(0.8),
             "Why Partner With Us?", 36, DBC_GREEN, True)

benefits = [
    ("New Foot Traffic", "Real customers walking through your door, incentivized by an exclusive offer."),
    ("Curated, Not Cheap", "This isn't a coupon app. Your bar is positioned as a members-only destination."),
    ("One-Time Offers", "Each member redeems once. They come back because they love your bar."),
    ("Zero Tech Hassle", "No hardware, no POS integration. Member shows a screen. That's it."),
    ("Data & Insights", "Track redemption volume, timing, and member activity at your venue."),
    ("Free to Join Pilot", "No cost, no commitment. See the results firsthand."),
]

for i, (title, desc) in enumerate(benefits):
    col = i % 2
    row = i // 2
    x = Inches(1 + col * 5.8)
    y = Inches(1.9 + row * 1.7)

    add_shape(slide, x, y + Inches(0.05), Inches(0.2), Inches(0.2), DBC_GREEN)
    add_text_box(slide, x + Inches(0.4), y - Inches(0.05), Inches(5), Inches(0.4),
                 title, 18, DARK_GRAY, True)
    add_text_box(slide, x + Inches(0.4), y + Inches(0.35), Inches(5), Inches(1),
                 desc, 14, MED_GRAY)


# ── SLIDE 4: Austin Pilot ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.15), DBC_GREEN)

add_text_box(slide, Inches(1), Inches(0.6), Inches(11), Inches(0.8),
             "Austin Pilot Program", 36, DBC_GREEN, True)

add_text_box(slide, Inches(1), Inches(1.5), Inches(10), Inches(0.8),
             "We're launching with a small group of hand-picked Austin bars. You choose the offer - a signature cocktail, a featured draft, a house special - whatever represents your bar.",
             20, DARK_GRAY)

pilot_details = [
    "Austin locals only (ZIP-verified membership)",
    "Free for bars during pilot - no cost, no commitment",
    "Full redemption tracking from day one",
    "You pick the drink offer that fits your brand",
    "Founding partners get priority placement as we scale",
]

add_bullet_list(slide, Inches(1), Inches(3.0), Inches(6), Inches(3.5), pilot_details, 18, DARK_GRAY)

# Callout box
box = add_shape(slide, Inches(8), Inches(3.0), Inches(4.3), Inches(2.8), LIGHT_BG)
add_text_box(slide, Inches(8.3), Inches(3.3), Inches(3.7), Inches(0.5),
             "PILOT PARTNERS GET", 14, DBC_GREEN, True)
pilot_perks = [
    "Early access & top placement",
    "Input on features & direction",
    "Founding partner recognition",
    "First look at growth data",
]
add_bullet_list(slide, Inches(8.3), Inches(3.9), Inches(3.7), Inches(2), pilot_perks, 15, MED_GRAY)


# ── SLIDE 5: Let's Partner Up ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DBC_DARK_GREEN)

add_text_box(slide, Inches(1.5), Inches(2.0), Inches(10), Inches(1.2),
             "Let's Partner Up", 48, WHITE, True, PP_ALIGN.CENTER)

add_shape(slide, Inches(5.5), Inches(3.3), Inches(2.3), Pt(2), WHITE)

add_text_box(slide, Inches(2), Inches(3.8), Inches(9), Inches(1),
             "We'd love to feature your bar in Dollar Bar Club.", 22, WHITE, False, PP_ALIGN.CENTER)

add_text_box(slide, Inches(2), Inches(4.8), Inches(9), Inches(1.5),
             "Get in touch to learn more and join the Austin pilot.\n\nBarGlance  |  dollarbarclub.com",
             20, WHITE, False, PP_ALIGN.CENTER)


# ── Save ──
output_path = r"C:\Users\jmoss\OneDrive\Desktop\Dollar Bar Club\Dollar Bar Club - Bar Partner Overview.pptx"
prs.save(output_path)
print(f"Saved to: {output_path}")
